from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "从零构建与预训练自回归语言模型"
    subtitle.text = "探索基础 Transformer 与现代大模型架构 (LLaMA-style) 优化\n汇报人: 宋奕洁 (2025231056)\nProject 2: Pre-training a GPT Model from Scratch"

    # Slide 2: Project Objectives
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "项目目标与核心挑战"
    tf = body_shape.text_frame
    tf.text = "核心目标："
    p = tf.add_paragraph()
    p.text = "彻底摆脱高度封装的库，使用原生 PyTorch 从零实现 GPT 模型。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "在自定义语料 (TinyStories/红楼梦) 上完成从分词、数据加载到自回归训练的全流程。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "高级挑战 (Advanced Tasks)："
    p.level = 0
    p = tf.add_paragraph()
    p.text = "架构深度优化：引入当前业界顶尖模型 (如 LLaMA, Qwen) 所采用的架构变体。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "推理与显存优化：实现 KV Cache 加速推理，并通过 Flash Attention 突破显存瓶颈。"
    p.level = 1

    # Slide 3: Baseline
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_shape.text = "基础架构实现 (Baseline Implementation)"
    tf = body_shape.text_frame
    tf.text = "模型骨架构建："
    p = tf.add_paragraph()
    p.text = "Embedding：Token Embedding 与绝对位置编码 (Absolute Positional Encoding)。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Causal Self-Attention：实现掩码多头注意力机制 (Masked MHA)，防止信息穿越。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "FeedForward Network：包含标准的 GELU 激活与全连接层。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "训练闭环：基于 AdamW 优化器，使用交叉熵损失函数优化 Next-Token Prediction。"
    p.level = 1

    # Slide 4: RoPE & RMSNorm
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_shape.text = "现代架构升级 —— RoPE 与 RMSNorm"
    tf = body_shape.text_frame
    tf.text = "旋转位置编码 (RoPE - Rotary Position Embedding)"
    p = tf.add_paragraph()
    p.text = "痛点：绝对位置编码在超出训练长度时泛化能力极差。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "改进：移除 nn.Embedding，在注意力层直接对 Query 和 Key 应用旋转矩阵。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "均方根归一化 (RMSNorm)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "痛点：传统的 LayerNorm 计算均值和方差，开销相对较大。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "改进：引入 RMSNorm 替换 LayerNorm，舍弃了均值中心化，提高计算速度。"
    p.level = 1

    # Slide 5: SwiGLU & GQA
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_shape.text = "现代架构升级 —— SwiGLU 与 GQA"
    tf = body_shape.text_frame
    tf.text = "门控激活网络 (SwiGLU)"
    p = tf.add_paragraph()
    p.text = "将原始的 GELU MLP 层替换为门控线性单元 SwiGLU。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "收益：极大增强了参数的非线性表达能力，获得更低的 Perplexity。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "分组查询注意力 (GQA - Grouped-Query Attention)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "痛点：多头注意力 (MHA) 在推理时 KV Cache 占用显存随长度线性暴增。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "改进：让多个 Query 共享同一组 Key 和 Value，显著降低显存占用。"
    p.level = 1

    # Slide 6: Flash Attention
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_shape.text = "极致性能 —— KV Cache 与 Flash Attention"
    tf = body_shape.text_frame
    tf.text = "KV Cache 机制"
    p = tf.add_paragraph()
    p.text = "在自回归生成时，缓存之前所有 Token 的 Key 和 Value。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "效果：时间复杂度从 O(N^2) 降为 O(N)，显著提升 Decode 速度。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Flash Attention 基准测试"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "引入了 PyTorch 原生的 scaled_dot_product_attention。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "效果：峰值显存占用大幅下降，计算耗时显著缩短。"
    p.level = 1

    # Slide 7: Generation
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_shape.text = "文本生成策略改进"
    tf = body_shape.text_frame
    tf.text = "在贪婪搜索之外，我们实现了复合的采样生成算法："
    p = tf.add_paragraph()
    p.text = "Temperature (温度系数)：调整 logits 的平滑度。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Top-K 采样：每次只从概率最高的 K 个词中采样，截断长尾低概率词。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Top-P (Nucleus Sampling)：基于累积概率阈值 P 进行动态词表截断。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "结合使用：生成既流畅又充满想象力的故事片段。"
    p.level = 1

    # Slide 8: Evaluation
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_shape.text = "评估与可视化 (Evaluation & Visualization)"
    tf = body_shape.text_frame
    tf.text = "定量评估 (Quantitative)"
    p = tf.add_paragraph()
    p.text = "不同规模模型 (Micro, Mini, Small) 的 Training Loss 和 Validation Perplexity 曲线。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "定性评估 (Qualitative)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "输入相同 Prompt，展示不同规模模型生成文本的连贯性和复杂度差异。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "注意力可视化 (Attention Heatmap)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "提取模型最后一步 Attention 权重，绘制词与词之间的注意力分布图。"
    p.level = 1

    # Slide 9: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_shape.text = "总结与未来展望 (Conclusion & Future Work)"
    tf = body_shape.text_frame
    tf.text = "项目总结"
    p = tf.add_paragraph()
    p.text = "成功从零构建了 GPT，并融合了 LLaMA 架构的核心组件 (RoPE, SwiGLU, RMSNorm, GQA)。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "通过 KV Cache 与 Flash Attention 实现了深度级别的性能调优。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "未来展望"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "通过 DPO (直接偏好优化) 让模型符合人类偏好。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "迁移至 DDP / FSDP，实现多卡分布式大规模预训练。"
    p.level = 1

    # Save
    prs.save('project2_presentation.pptx')
    print("Presentation saved successfully.")

if __name__ == '__main__':
    create_presentation()
